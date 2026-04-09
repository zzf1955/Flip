"""生成汇报 PPT：第一人称人形机器人视频生成"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── 颜色方案（简洁学术风）──
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x1A, 0x1A, 0x2E)
DARK_BLUE = RGBColor(0x1E, 0x3A, 0x5F)
MID_BLUE = RGBColor(0x2C, 0x5F, 0x8A)
LIGHT_BLUE = RGBColor(0xE8, 0xF0, 0xF8)
ACCENT = RGBColor(0xC0, 0x39, 0x2B)
GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xAA, 0xAA, 0xAA)
GREEN = RGBColor(0x27, 0xAE, 0x60)
TABLE_HEADER_BG = RGBColor(0x1E, 0x3A, 0x5F)
TABLE_ALT_BG = RGBColor(0xF5, 0xF8, 0xFC)
TABLE_WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height


# ── 工具函数 ──
def add_bg(slide, color=WHITE):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = color


def add_top_bar(slide):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.08)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = MID_BLUE
    shape.line.fill.background()


def add_page_num(slide, num, total):
    txBox = slide.shapes.add_textbox(
        Inches(12.3), Inches(7.05), Inches(0.8), Inches(0.35)
    )
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"{num}/{total}"
    p.font.size = Pt(10)
    p.font.color.rgb = LIGHT_GRAY
    p.alignment = PP_ALIGN.RIGHT


def title_slide(title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_bg(slide, DARK_BLUE)
    # 左侧竖条
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.5), Inches(0.06), Inches(3.5)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    # title
    txBox = slide.shapes.add_textbox(Inches(1.2), Inches(1.8), Inches(10), Inches(2.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(38)
    p.font.color.rgb = WHITE
    p.font.bold = True
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(18)
        p2.font.color.rgb = RGBColor(0xBB, 0xCC, 0xDD)
        p2.space_before = Pt(20)
    return slide


def section_slide(title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, MID_BLUE)
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT
    # 下划线装饰
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1), Inches(4.2), Inches(3), Inches(0.05)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = WHITE
    line.line.fill.background()
    return slide


def content_slide(title, bullets, sub_bullets=None):
    """普通内容页。bullets 是列表，sub_bullets 是 dict: {idx: [子项]}"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_top_bar(slide)
    # title
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(26)
    p.font.color.rgb = DARK_BLUE
    p.font.bold = True
    # body
    txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(11.5), Inches(5.5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    sub_bullets = sub_bullets or {}
    for i, b in enumerate(bullets):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p.text = b
        p.font.size = Pt(18)
        p.font.color.rgb = BLACK
        p.space_before = Pt(8)
        p.level = 0
        if i in sub_bullets:
            for sb in sub_bullets[i]:
                sp = tf2.add_paragraph()
                sp.text = sb
                sp.font.size = Pt(15)
                sp.font.color.rgb = GRAY
                sp.space_before = Pt(4)
                sp.level = 1
    return slide


def two_col_slide(title, left_title, left_bullets, right_title, right_bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_top_bar(slide)
    # title
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(26)
    p.font.color.rgb = DARK_BLUE
    p.font.bold = True

    for col_x, col_title, col_bullets in [
        (0.8, left_title, left_bullets),
        (6.8, right_title, right_bullets),
    ]:
        # 小标题
        txB = slide.shapes.add_textbox(
            Inches(col_x), Inches(1.3), Inches(5.5), Inches(0.5)
        )
        tff = txB.text_frame
        pp = tff.paragraphs[0]
        pp.text = col_title
        pp.font.size = Pt(20)
        pp.font.color.rgb = MID_BLUE
        pp.font.bold = True
        # 分割线
        ln = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(col_x),
            Inches(1.85),
            Inches(5),
            Inches(0.025),
        )
        ln.fill.solid()
        ln.fill.fore_color.rgb = MID_BLUE
        ln.line.fill.background()
        # bullets
        txB2 = slide.shapes.add_textbox(
            Inches(col_x), Inches(2.0), Inches(5.5), Inches(5.0)
        )
        tff2 = txB2.text_frame
        tff2.word_wrap = True
        for j, b in enumerate(col_bullets):
            pp2 = tff2.paragraphs[0] if j == 0 else tff2.add_paragraph()
            pp2.text = b
            pp2.font.size = Pt(16)
            pp2.font.color.rgb = BLACK
            pp2.space_before = Pt(6)
    return slide


def add_table_slide(title, headers, rows, col_widths=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_top_bar(slide)
    # title
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(26)
    p.font.color.rgb = DARK_BLUE
    p.font.bold = True

    n_rows = len(rows) + 1
    n_cols = len(headers)
    if col_widths is None:
        total_w = 11.5
        col_widths = [total_w / n_cols] * n_cols

    table_shape = slide.shapes.add_table(
        n_rows,
        n_cols,
        Inches(0.8),
        Inches(1.5),
        sum(Inches(w) for w in col_widths),
        Inches(0.45 * n_rows),
    )
    table = table_shape.table

    for c_i in range(n_cols):
        table.columns[c_i].width = Inches(col_widths[c_i])

    # header
    for c_i, h in enumerate(headers):
        cell = table.cell(0, c_i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = TABLE_HEADER_BG
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(14)
            paragraph.font.color.rgb = WHITE
            paragraph.font.bold = True
            paragraph.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # rows
    for r_i, row in enumerate(rows):
        for c_i, val in enumerate(row):
            cell = table.cell(r_i + 1, c_i)
            cell.text = str(val)
            if r_i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ALT_BG
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_WHITE
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(13)
                paragraph.font.color.rgb = BLACK
                paragraph.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    return slide


# ══════════════════════════════════════════════
#  开始生成幻灯片
# ══════════════════════════════════════════════
TOTAL = 20

# ── 1. 封面 ──
title_slide(
    "第一人称人形机器人视频生成",
    "反向数据构造：真实机器人视频上合成人手",
)

# ── 2. 目录 ──
content_slide(
    "目录",
    [
        "一、上次讨论回顾与新文献",
        "二、目前思路的调整",
        "三、问题背景：Human-to-Robot 视频生成现状",
        "四、方法：反向数据构造 — Robot → Human",
        "五、目前进度",
        "六、待解决的问题",
    ],
)

# ── 3. Section: 上次讨论回顾 ──
section_slide("一、上次讨论回顾与新文献")

# ── 4. 两个方向总览 ──
two_col_slide(
    "上次提到的两个方向",
    "方向一：Latent 空间动作表示（DisMo）",
    [
        "• Stability AI, NeurIPS 2025",
        "• 用 latent motion 表征绕开像素空间",
        "• 训练数据：280 万视频片段（4.9k 小时）",
        "• 86M 参数 3D ViT + 675M DiT-XL",
        "",
        "复现结论：",
        "• 表示的是整个视频的动态",
        "• 只能表示简单动作，精细人体动作差",
        "• 仍然绕不开像素空间编辑",
    ],
    "方向二：外观替换（MIMO）",
    [
        "• 阿里通义实验室",
        "• 深度分层空间分解建模",
        "• 训练数据：HUD-7K（5K 真实 + 2K 合成视频）",
        "• 基于 SD 1.5 + AnimateDiff",
        "",
        "复现结论：",
        "• 依赖 SMPL 人体模型，机器人无法拟合",
        "• Detectron2 对机器人检测置信度低",
        "• 肢体遮挡场景频繁失败",
    ],
)

# ── 5. Cosmos Transfer ──
content_slide(
    "之前的想法：Cosmos Transfer",
    [
        "缺乏强条件约束，参考帧并不是强制首帧",
        "首帧注入会有漂移，效果非常差",
        "结论：不可行",
    ],
)

# ── 6. Section: 思路调整 ──
section_slide("二、目前思路的调整")

# ── 7. 第三→第一人称 ──
content_slide(
    "第三人称 → 第一人称",
    [
        "放弃第三人称的三个原因：",
        "",
        "1. 数据太少",
        "2. 遮挡难处理",
        "3. 机器人识别度低",
        "",
        "第一人称的优势：",
        "",
        "• 宇树 G1 第一人称数据自带姿态数据（关节编码器）",
        "• 可结合 MIMO 的做法做精准渲染",
        "• Ego 视角数据量相对充足",
    ],
)

# ── 8. 编辑方法的问题 ──
two_col_slide(
    "编辑方法的问题与结论",
    "生成式编辑的局限",
    [
        "• Ego 数据多为鱼眼透视",
        "• 文生图模型对鱼眼 OOD，效果差",
        "• ControlNet 在边缘情况下效果差",
        "  （可能因为卷积对畸变敏感）",
    ],
    "结论",
    [
        "• 回到第一人称视角",
        "• 使用更传统的编辑方法",
        "• 利用 G1 自带的姿态数据",
        "  做精准渲染，而非纯生成式编辑",
    ],
)

# ── 9. Section: 问题背景 ──
section_slide("三、问题背景")

# ── 10. 维度一：视角 × 机器人类型 ──
add_table_slide(
    "维度一：视角 × 机器人类型",
    ["", "机械臂", "人形机器人"],
    [
        ["第一人称", "Masquerade, Mitty", "空白（本工作）"],
        ["第三人称", "Human2Robot", "X-Humanoid, Dream2Act"],
    ],
    col_widths=[2.5, 4.5, 4.5],
)

# ── 11. 数据量对比（从论文中查到的） ──
add_table_slide(
    "各工作数据量对比",
    ["工作", "视角", "机器人", "训练数据", "目标端"],
    [
        [
            "Masquerade",
            "第一人称",
            "机械臂",
            "675K 帧 (Epic Kitchens)\n+ 50 真实 demo",
            "渲染合成",
        ],
        [
            "Mitty",
            "第一人称",
            "机械臂",
            "6,000 合成 + 1,019 真实配对",
            "渲染合成",
        ],
        [
            "X-Humanoid",
            "第三人称",
            "人形 (Optimus)",
            "11,172 UE5 配对 (17+ 小时)",
            "渲染合成（两端）",
        ],
        [
            "Human2Robot",
            "第三人称",
            "机械臂",
            "2,600 VR 遥操作 episodes",
            "真实",
        ],
        [
            "Dream2Act",
            "第三人称",
            "人形 (G1)",
            "零样本（无训练数据）",
            "零样本生成",
        ],
        [
            "本工作",
            "第一人称",
            "人形 (G1)",
            "真实 G1 数据 + 反向构造",
            "真实",
        ],
    ],
    col_widths=[1.8, 1.5, 2.0, 3.5, 2.5],
)

# ── 12. 维度二：数据构造方向 ──
content_slide(
    "维度二：数据构造方向",
    [
        "现有工作均为 Human → Robot 方向：",
        "",
        "• Masquerade / Mitty：人类视频 → 去手 → inpaint → 渲染机械臂",
        "      得到 (真 human, 合成 robot)",
        "",
        "• X-Humanoid：UE5 渲染配对 human/humanoid 动画",
        "      得到 (合成 human, 合成 robot)",
        "",
        "共同问题：目标端（robot 视频）是渲染/合成的，存在 sim2real gap",
    ],
)

# ── 13. Masquerade 详述 ──
content_slide(
    "Masquerade (ICLR 2026, Stanford)",
    [
        "任务：第一人称人类视频 → 机械臂视频（rendering overlay）",
        "",
        "Pipeline：HaMeR 提手部 pose → Detectron2+SAM2 去手 → E2FGVI inpaint → 渲染机械臂",
        "",
        "数据：675K 帧 (Epic Kitchens) + 50 真实 bimanual robot demo",
        "",
        "结果：OOD 场景平均成功率 12% → 74%（提升约 5-6×）",
        "",
        "训练：ViT-Base 150K steps + Diffusion Policy 40K steps",
    ],
    {
        7: [
            "局限：rule-based overlay，累积误差大；仅限特定机械臂；目标是渲染的",
            "手部 pose 估计在快速运动和遮挡下失败率高",
        ]
    },
)

# ── 14. Mitty 详述 ──
content_slide(
    "Mitty (ICLR 2026 under review)",
    [
        "任务：第一人称 human → robot 端到端视频生成（diffusion）",
        "",
        "方法：Wan 2.2 + LoRA (rank 96)，双向注意力 video in-context learning",
        "",
        "数据：6,000 合成配对 (Masquerade pipeline) + 1,019 真实配对 (H2R)",
        "",
        "结果：任务成功率 93.7%（vs Masquerade 31.2%）",
        "",
        "训练：TI2V-14B 5K steps / TI2V-5B 20K steps，2× H200",
    ],
    {
        7: [
            "局限：仅限机械臂，不支持人形机器人",
            "只能生成视频，不能输出 action 序列",
            "依赖文本描述；不保证机械臂物理结构",
        ]
    },
)

# ── 15. X-Humanoid & Dream2Act ──
two_col_slide(
    "X-Humanoid & Dream2Act",
    "X-Humanoid (2512)",
    [
        "• 第三人称 human → humanoid (Tesla Optimus)",
        "• Wan 2.2 T2V-5B + LoRA (rank 96)",
        "• 11,172 UE5 配对片段 (17+ 小时)",
        "• 仅 500 steps，2.5 小时训练（4× H200）",
        "• 用户研究：运动一致性 69%",
        "",
        "局限：",
        "• 仅限第三人称",
        "• 合成数据两端都是渲染的",
        "• 超过 500 steps 真实视频效果下降",
    ],
    "Dream2Act (2603)",
    [
        "• 零样本人形 (G1) 视频'幻想'",
        "• Seedance 2.0 图片+文本 → 机器人视频",
        "• 提取关节轨迹执行",
        "• 20 trials/task，整体成功率 37.5%",
        "  - 踢球 40%, 拥抱 10%, 击打 30%, 坐下 70%",
        "",
        "局限：",
        "• 第三人称",
        "• 零样本精度有限",
        "• 仅粗粒度全身动作",
    ],
)

# ── 16. Section: 方法 ──
section_slide("四、方法：反向数据构造")

# ── 17. 核心观察 ──
content_slide(
    "核心观察：反转数据构造方向",
    [
        "现有方向：真 human → 渲染 robot → (真 human, 合成 robot) → 训练模型",
        "                                          ↑ sim2real gap 在目标端",
        "",
        "本工作：  真 robot → 合成 human → (合成 human, 真 robot) → 训练模型",
        "                                          ↑ domain gap 在输入端（影响更小）",
        "",
        "为什么 domain gap 在输入端更好？",
        "",
        "1. 目标端是真实的 → 模型学到的是生成真实机器人视频的分布",
        "2. 推理时输入是真实人类视频，质量 ≥ 合成 → 泛化方向「从差到好」",
        "3. Masquerade 已证明 imperfect overlay 也能提升 5-6× → 输入端容忍度很高",
    ],
)

# ── 18. Pipeline ──
content_slide(
    "Pipeline 设计",
    [
        "G1 第一人称视频 + G1 姿态数据（关节编码器）",
        "",
        "Step 1  Pose 获取：直接读取 G1 关节角度（29 DOF）",
        "            → 精确的本体感知，无需视觉 pose estimation（关键优势）",
        "",
        "Step 2  Robot 手臂分割+去除：Pinocchio FK → Mesh 投影 → 精确 Mask",
        "            → E2FGVI 背景 inpainting",
        "",
        "Step 3  运动学映射：Robot joint state → Human pose",
        "",
        "Step 4  Human 手部渲染：在 inpainted 背景上渲染人手",
        "",
        "Step 5  得到 (合成 human, 真实 robot) 配对 → 微调 Wan 2.2 + LoRA",
        "",
        "（可选）结合 World Action Model 范式，联合预测 action + video",
    ],
)

# ── 19. 优势 ──
content_slide(
    "方法优势",
    [
        "1. Pose 获取精度高",
        "",
        "2. 目标是真实视频",
        "",
        "3. 仅需机器人数据",
        "",
        "4. 可结合 World Action Model",
    ],
    {
        0: ["用本体感知而非视觉估计，消除 Masquerade 最大误差源"],
        2: ["无 sim2real gap，模型直接学习真实机器人视频分布"],
        4: ["不需手工构造配对视频，只需更换机器人数据源，pipeline 可复用"],
        6: ["联合预测 action 和 video 帧，提升物理精准性"],
    },
)

# ── 20. Section: 进度 ──
section_slide("五、目前进度")

# ── 21. 进度详情 ──
content_slide(
    "目前进度",
    [
        "✓ 姿态数据获取",
        "",
        "✓ 精准渲染 Mask",
        "",
        "    已完成工作：",
        "    • Pinocchio 正向运动学 + STL Mesh 加载",
        "    • OpenCV 鱼眼相机模型投影",
        "    • PSO 自动标定（14 参数），IoU ≈ 0.897",
        "    • 多任务、多 episode 批量 overlay 验证",
        "",
        "（效果图和视频见演示）",
    ],
)

# ── 22. Section: 待解决 ──
section_slide("六、待解决的问题")

# ── 23. 待解决问题 ──
content_slide(
    "待解决的问题",
    [
        "1. Human hand rendering 质量",
        "",
        "2. 运动学映射 G1 → Human",
        "",
        "3. Mask 擦除效果",
    ],
    {
        0: [
            "渲染人手的视觉质量如何？",
            "视频连贯性如何保证？",
        ],
        2: [
            "G1 关节 → 人体关节的映射是否自然？",
            "自由度差异如何处理？",
        ],
        4: [
            "目前 mask 覆盖率较低（IoU ~0.9，仍有遗漏区域）",
            "inpainting 对残留 artifact 的容忍度未知",
        ],
    },
)

# ── 24. 谢谢 ──
title_slide("谢谢", "")

# ── 添加页码 ──
for i, slide in enumerate(prs.slides):
    add_page_num(slide, i + 1, len(prs.slides))

# ── 保存 ──
out_path = "/Users/zzf/share/VideoEdit/report.pptx"
prs.save(out_path)
print(f"PPT saved to {out_path}")
print(f"Total slides: {len(prs.slides)}")
